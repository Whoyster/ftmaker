from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os
import json
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_shape_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))

def generate_ppt(id_title:str="id_title", contents_dir:str="results"):

    with open(contents_dir+'/'+id_title+'/'+id_title+'.json') as ifs:
        contents = json.load(ifs)
        print(contents)

    # 새 프레젠테이션 생성
    prs = Presentation()

    # A4 용지 크기로 변경
    #prs.slide_width = Inches(11.69)
    #prs.slide_height = Inches(8.27)

    # 제목 슬라이드 추가
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = contents['title']

    for page in contents["pages"]:
        img_path = page['img_path']
        text_ko = page['txt']
        add_scene(prs, img_path, text_ko)

    # PPT 파일 저장
    prs.save(contents_dir+'/'+id_title+'/'+id_title+'.pptx')
    print('saved '+contents_dir+'/'+id_title+'/'+id_title+'.pptx')

def add_scene(prs, img_path: str, text_ko: str):
    # add empty slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # add image full size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    pic = slide.shapes.add_picture(img_path, left=0, top=0, width=slide_width, height=slide_height)

    # add white rounded rectangle bottom centered
    text_box_width = prs.slide_width * 0.9
    text_box_height = prs.slide_height * 0.15
    text_box_left = (slide_width - text_box_width) / 2
    text_box_top = slide_height - text_box_height - Inches(0.5)
    text_box = slide.shapes.add_shape(
        5,  # rounded rectangle
        text_box_left, text_box_top, text_box_width, text_box_height
    )
    text_box.adjustment_values = (0.15, 0.15, 0.15, 0.15)

    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    _set_shape_transparency(text_box,80000)
    
    line = text_box.line
    line.color.rgb = RGBColor(255, 255, 255)
    line.width = Pt(0)
    line.dash_style = None

    tf = text_box.text_frame
    tf.text = text_ko
    for p in tf.paragraphs:
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = 'NanumBarunGothic'
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = 1  # centered
        p.line_spacing = 1.5
        p.space_before = 0
        p.space_after = 0


